"""
Generate Transaction Intelligence Presentation.pptx
Comprehensive executive + technical PowerPoint for TigerGraph hybrid search demo.

Run: python scripts/generate_presentation.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

ROOT   = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "Transaction_Intelligence_Presentation.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x10, 0x1B, 0x2E)   # slide background / dark elements
BLUE    = RGBColor(0x00, 0x73, 0xCF)   # TigerGraph blue accent
ORANGE  = RGBColor(0xF1, 0x5A, 0x22)   # TigerGraph orange accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE8, 0xF0, 0xFE)   # soft blue-white for backgrounds
SILVER  = RGBColor(0xB0, 0xBE, 0xCE)   # muted text
TEAL    = RGBColor(0x00, 0xAF, 0xD8)   # highlight colour
GREEN   = RGBColor(0x2A, 0x9D, 0x5C)
AMBER   = RGBColor(0xF5, 0xA6, 0x23)
DARKBG  = RGBColor(0x1A, 0x25, 0x3A)   # content-slide background stripe

# Slide dimensions — 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Low-level helpers ────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_width: int = 0) -> object:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width     = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text: str, left, top, width, height,
                 font_size: int = 14, bold: bool = False, italic: bool = False,
                 color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
                 wrap: bool = True, font_name: str = "Calibri") -> object:
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text      = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_para(tf, text: str, font_size: int = 13, bold: bool = False,
             color: RGBColor = WHITE, indent: int = 0,
             space_before: int = 0, align=PP_ALIGN.LEFT,
             font_name: str = "Calibri") -> None:
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if indent:
        p.level = indent
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name


def slide_header(slide, title: str, subtitle: str = "",
                 bar_color: RGBColor = BLUE) -> None:
    """Dark header bar spanning full width."""
    add_rect(slide, 0, 0, 13.33, 1.05, fill=NAVY)
    add_rect(slide, 0, 1.0, 13.33, 0.08, fill=bar_color)
    add_text_box(slide, title, 0.35, 0.12, 11.5, 0.65,
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.35, 0.68, 11.5, 0.35,
                     font_size=13, color=SILVER, align=PP_ALIGN.LEFT)


def bullet_box(slide, items: list[str], left, top, width, height,
               font_size: int = 13, color: RGBColor = WHITE,
               title: str = "", title_color: RGBColor = TEAL,
               bg: RGBColor | None = None, icon: str = "•") -> None:
    if bg:
        add_rect(slide, left, top, width, height, fill=bg)
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"
        p.space_before = Pt(2)
        first = False
    for item in items:
        p = tf.add_paragraph() if not first else (tf.paragraphs[0] if first and not title else tf.add_paragraph())
        first = False
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_before = Pt(3)


def card(slide, left, top, width, height, title: str, body: str,
         bg: RGBColor = DARKBG, title_color: RGBColor = TEAL,
         body_color: RGBColor = WHITE, accent: RGBColor = BLUE) -> None:
    add_rect(slide, left, top, width, height, fill=bg)
    add_rect(slide, left, top, 0.06, height, fill=accent)
    add_text_box(slide, title, left + 0.15, top + 0.12, width - 0.2, 0.3,
                 font_size=13, bold=True, color=title_color)
    add_text_box(slide, body, left + 0.15, top + 0.42, width - 0.2, height - 0.55,
                 font_size=11, color=body_color, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 0, 0.18, 7.5, fill=ORANGE)
add_rect(sl, 0.18, 3.5, 13.15, 0.07, fill=BLUE)

add_text_box(sl, "TRANSACTION INTELLIGENCE", 0.5, 1.5, 12.5, 0.9,
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(sl, "Powered by TigerGraph 4.2  ·  Vector + Graph Hybrid Search",
             0.5, 2.55, 12.0, 0.55, font_size=20, color=TEAL, align=PP_ALIGN.LEFT)
add_text_box(sl, "A unified intelligence platform for transaction risk detection, "
             "behavioural analytics, and semantic search — built on native graph + vector technology.",
             0.5, 3.25, 10.5, 0.8, font_size=14, color=SILVER, align=PP_ALIGN.LEFT)

add_text_box(sl, "5,000 Transactions  ·  220 Merchants  ·  848 Users  ·  384-dim Embeddings",
             0.5, 4.3, 12.0, 0.4, font_size=13, color=AMBER, align=PP_ALIGN.LEFT)

add_text_box(sl, "⚡ Hybrid Search  |  🔬 Cluster Explorer  |  🧠 GraphSAGE",
             0.5, 4.9, 10.0, 0.4, font_size=13, color=SILVER, align=PP_ALIGN.LEFT)

add_text_box(sl, "CONFIDENTIAL — DEMO BUILD", 0.5, 6.9, 5.0, 0.35,
             font_size=9, color=RGBColor(0x55, 0x65, 0x75), align=PP_ALIGN.LEFT)
add_text_box(sl, "TigerGraph 4.2.2  |  sentence-transformers  |  HNSW Index",
             7.5, 6.9, 5.5, 0.35, font_size=9, color=RGBColor(0x55, 0x65, 0x75),
             align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE CHALLENGE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "The Challenge with Traditional Transaction Analytics",
             "Why rule engines and SQL alone are no longer enough", bar_color=ORANGE)

problems = [
    ("Rule Engine Limits",
     "Rules are static and require manual updating. Every new fraud pattern\n"
     "needs a new rule — and adversaries adapt faster than analysts can write them.",
     ORANGE),
    ("SQL / BI Query Blindspots",
     "Aggregate queries miss relationships between entities. A fraudster operating\n"
     "across 12 merchants and 300 users is invisible to column-level analysis.",
     BLUE),
    ("No Semantic Understanding",
     "Keyword and category searches cannot find conceptually similar transactions.\n"
     "\"Suspicious late-night cash-like activity\" has no SQL equivalent.",
     TEAL),
    ("Cold-Start Problem",
     "A brand new merchant or user has no history to flag against. Rule engines\n"
     "are blind until the first confirmed fraud event — which is already too late.",
     GREEN),
]

cols = [(0.4, 1.35), (3.55, 1.35), (6.7, 1.35), (9.85, 1.35)]
for (left, top), (title, body, col) in zip(cols, problems):
    add_rect(sl, left, top, 2.9, 4.5, fill=NAVY)
    add_rect(sl, left, top, 2.9, 0.06, fill=col)
    add_text_box(sl, title, left + 0.15, top + 0.2, 2.6, 0.45,
                 font_size=13, bold=True, color=col)
    add_text_box(sl, body, left + 0.15, top + 0.75, 2.6, 3.5,
                 font_size=11, color=WHITE, wrap=True)

add_rect(sl, 0.4, 6.0, 12.53, 0.6, fill=RGBColor(0xEB, 0xF5, 0xFF))
add_text_box(sl, "📌  Transaction Intelligence addresses all four gaps simultaneously — "
             "using TigerGraph's native vector index for semantic search and graph traversal "
             "for relationship-aware context expansion.",
             0.6, 6.05, 12.2, 0.5, font_size=12, color=NAVY, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PLATFORM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Platform Architecture",
             "End-to-end flow from raw transaction data to intelligent search results")

# Layer labels (left spine)
layers = [
    (1.15, "1  DATA INGESTION",      ORANGE),
    (2.15, "2  FEATURE ENGINEERING", BLUE),
    (3.15, "3  EMBEDDING LAYER",     TEAL),
    (4.4,  "4  TIGERGRAPH CORE",     RGBColor(0xA0, 0x60, 0xFF)),
    (5.7,  "5  QUERY ENGINE",        GREEN),
    (6.6,  "6  PRESENTATION",        AMBER),
]
for top, label, col in layers:
    add_rect(sl, 0.15, top, 0.12, 0.55, fill=col)
    add_text_box(sl, label, 0.35, top + 0.08, 1.7, 0.4,
                 font_size=8, bold=True, color=col)

# ── Row 1: Data Sources ──────────────────────────────────────────────────────
sources = ["5,000\nTransactions CSV", "220\nMerchants CSV", "848\nUsers CSV",
           "Edge\nRelationships"]
sx = [2.3, 3.9, 5.5, 7.1]
for x, s in zip(sx, sources):
    add_rect(sl, x, 1.1, 1.4, 0.65, fill=RGBColor(0x1E, 0x35, 0x55), line=ORANGE, line_width=1)
    add_text_box(sl, s, x + 0.05, 1.15, 1.3, 0.55, font_size=9, color=WHITE,
                 align=PP_ALIGN.CENTER)

add_rect(sl, 8.7, 1.1, 2.0, 0.65, fill=RGBColor(0x1E, 0x35, 0x55), line=ORANGE, line_width=1)
add_text_box(sl, "TigerGraph\nSchema (GSQL)", 8.75, 1.15, 1.9, 0.55,
             font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

# ── Row 2: Feature Engineering ───────────────────────────────────────────────
feats = ["Amount Band\nLabelling", "MCC Text\nMapping", "Time Semantics\n(hour / weekday)",
         "Risk Tag\nDerivation", "Merchant Category\nText"]
fx = [2.3, 3.75, 5.2, 6.65, 8.1]
for x, f in zip(fx, feats):
    add_rect(sl, x, 2.1, 1.3, 0.65, fill=RGBColor(0x10, 0x2A, 0x48), line=BLUE, line_width=1)
    add_text_box(sl, f, x + 0.05, 2.15, 1.2, 0.55, font_size=9, color=WHITE,
                 align=PP_ALIGN.CENTER)

# ── Row 3: Embeddings ────────────────────────────────────────────────────────
embs = [
    ("Txn Risk\nEmbedding\n384-dim", ORANGE),
    ("Txn Behaviour\nEmbedding\n384-dim", BLUE),
    ("Merchant\nEmbedding\n384-dim", TEAL),
    ("User\nEmbedding\n384-dim", GREEN),
]
ex = [2.3, 4.15, 6.0, 7.85]
for x, (label, col) in zip(ex, embs):
    add_rect(sl, x, 3.1, 1.7, 0.8, fill=RGBColor(0x0A, 0x20, 0x35), line=col, line_width=2)
    add_text_box(sl, label, x + 0.05, 3.15, 1.6, 0.7, font_size=9, color=col,
                 align=PP_ALIGN.CENTER, bold=True)
add_rect(sl, 9.7, 3.1, 2.3, 0.8, fill=RGBColor(0x0A, 0x20, 0x35),
         line=RGBColor(0xA0, 0x60, 0xFF), line_width=1)
add_text_box(sl, "sentence-transformers\nall-MiniLM-L6-v2", 9.75, 3.15, 2.2, 0.7,
             font_size=9, color=RGBColor(0xC0, 0x90, 0xFF), align=PP_ALIGN.CENTER)

# ── Row 4: TigerGraph Core ───────────────────────────────────────────────────
add_rect(sl, 2.1, 4.1, 9.0, 1.1, fill=RGBColor(0x15, 0x25, 0x45))
add_rect(sl, 2.1, 4.1, 9.0, 0.06, fill=RGBColor(0xA0, 0x60, 0xFF))
add_text_box(sl, "TigerGraph 4.2.2 Cloud", 2.2, 4.12, 3.0, 0.35,
             font_size=11, bold=True, color=RGBColor(0xC0, 0x90, 0xFF))
tg_nodes = [
    ("Graph Store\n6 Vertex Types\n5 Edge Types", RGBColor(0xA0, 0x60, 0xFF)),
    ("HNSW Vector Index\nCOSINE · 384-dim\nFloat32", TEAL),
    ("GSQL Engine\nContext Queries\nSummary Queries", BLUE),
    ("REST++ API\nVector Search\nInstalled Queries", ORANGE),
]
tx = [2.2, 4.4, 6.6, 8.8]
for x, (label, col) in zip(tx, tg_nodes):
    add_rect(sl, x, 4.25, 2.0, 0.8, fill=RGBColor(0x0D, 0x1D, 0x35), line=col, line_width=1)
    add_text_box(sl, label, x + 0.07, 4.3, 1.85, 0.7, font_size=9,
                 color=col, align=PP_ALIGN.CENTER)

# ── Row 5: Query Engine ───────────────────────────────────────────────────────
qnodes = [
    ("NL Query\n↓ Embed", WHITE),
    ("Vector\nTop-K Search", TEAL),
    ("Graph\nTraversal", BLUE),
    ("Context\nExpansion", GREEN),
    ("Summary\n& Results", AMBER),
]
qx = [2.3, 4.0, 5.7, 7.4, 9.1]
for x, (label, col) in zip(qx, qnodes):
    add_rect(sl, x, 5.6, 1.55, 0.7, fill=RGBColor(0x10, 0x20, 0x30), line=col, line_width=1)
    add_text_box(sl, label, x + 0.05, 5.65, 1.45, 0.6, font_size=9,
                 color=col, align=PP_ALIGN.CENTER)
    if x < 9.1:
        add_text_box(sl, "→", x + 1.55, 5.88, 0.2, 0.25, font_size=12, color=SILVER)

# ── Row 6: Presentation layer ─────────────────────────────────────────────────
apps = ["🔍 Hybrid Search\n(4 modes)", "🔬 Cluster Explorer\n(UMAP 2D)", "🧠 GraphSAGE\n(1-hop / 2-hop)"]
ax = [2.5, 5.5, 8.5]
for x, label in zip(ax, apps):
    add_rect(sl, x, 6.55, 2.6, 0.65, fill=RGBColor(0x1A, 0x2E, 0x22), line=AMBER, line_width=1)
    add_text_box(sl, label, x + 0.1, 6.6, 2.4, 0.55, font_size=10,
                 color=AMBER, align=PP_ALIGN.CENTER)

# Vertical arrows between layers
for top in [1.77, 2.77, 3.92, 5.22, 6.32]:
    add_text_box(sl, "▼", 6.3, top, 0.4, 0.25, font_size=12, color=SILVER, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GRAPH DATA MODEL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "Graph Data Model  —  Tran_graph",
             "6 vertex types · 5 directed edge types · vector embeddings on 3 entities", bar_color=BLUE)

vertex_specs = [
    ("Transaction",      ORANGE, 5.5, 2.3,
     ["Transaction_id (PK)", "Amount : STRING", "amount_float : FLOAT",
      "Channel_type : STRING", "Currency_code : UINT", "response_code : STRING",
      "▸ risk_embedding : VECTOR[384]", "▸ behaviour_emb : VECTOR[384]"]),
    ("Merchant",         TEAL,   1.5, 2.3,
     ["merchantid (PK)", "merch_name : STRING", "merch_type : STRING",
      "acquirer_id : UINT", "ica_code : INT",
      "▸ embedding : VECTOR[384]"]),
    ("USER",             BLUE,   9.5, 2.3,
     ["userid (PK)", "dw_issuer_id : UINT", "dw_product_cd : UINT",
      "country_cd : FLOAT", "expiry_date : FLOAT",
      "▸ embedding : VECTOR[384]"]),
    ("Location",         GREEN,  1.5, 5.0,
     ["locationid (PK)", "merch_city : STRING",
      "merch_region_code : STRING", "merchant_country_code : STRING"]),
    ("Transaction_DateTime", AMBER, 5.5, 5.0,
     ["PanID : STRING", "TransactionDate : STRING", "TransactionTime : STRING"]),
    ("MCC_Code",         RGBColor(0xD0, 0x60, 0xC0), 9.5, 5.0,
     ["mcc_code : STRING"]),
]

for (vtype, col, lft, top, attrs) in vertex_specs:
    bw = 2.6 if vtype not in ("Location", "MCC_Code") else 2.4
    bh = 2.3 if top < 4 else 1.8
    add_rect(sl, lft, top, bw, bh, fill=NAVY, line=col, line_width=2)
    add_text_box(sl, vtype, lft + 0.1, top + 0.1, bw - 0.2, 0.35,
                 font_size=12, bold=True, color=col)
    add_rect(sl, lft, top + 0.45, bw, 0.02, fill=col)
    y = top + 0.55
    for attr in attrs:
        bold = attr.startswith("▸")
        c = TEAL if bold else SILVER
        add_text_box(sl, attr, lft + 0.1, y, bw - 0.2, 0.26,
                     font_size=9, color=c, bold=bold)
        y += 0.24

# Edges as connector labels
edges = [
    (3.95, 3.25, "Initiate"),
    (7.95, 3.25, "Occurs"),
    (5.5,  4.55, "Happens_at"),
    (3.0,  4.55, "Has"),
    (8.5,  4.55, "Categorized_by"),
]
for lx, ty, elabel in edges:
    add_rect(sl, lx, ty, 1.6, 0.35, fill=RGBColor(0x2A, 0x35, 0x50), line=WHITE, line_width=0)
    add_text_box(sl, f"── {elabel} ──", lx + 0.05, ty + 0.04, 1.5, 0.28,
                 font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, 0.3, 6.1, 12.7, 0.5, fill=RGBColor(0xEB, 0xF5, 0xFF))
add_text_box(sl, "▸ Vector attributes:  FLOAT · 384 dimensions · COSINE metric · HNSW index type  "
             "—  configured identically across Transaction (×2), Merchant (×1) and USER (×1).",
             0.5, 6.15, 12.3, 0.4, font_size=11, color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TIGERGRAPH VECTOR CAPABILITIES
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "TigerGraph Native Vector Capabilities",
             "Why TigerGraph 4.x is uniquely positioned for hybrid search", bar_color=ORANGE)

caps = [
    ("HNSW Index",
     "Hierarchical Navigable Small World graph index — industry-standard\n"
     "approximate nearest neighbour algorithm. Sub-millisecond Top-K\n"
     "retrieval across millions of vectors.",
     ORANGE, 0.4, 1.3),
    ("COSINE Similarity",
     "Measures directional similarity between embedding vectors.\n"
     "Ideal for semantic text embeddings where magnitude is normalised\n"
     "and angle captures meaning.",
     TEAL, 4.55, 1.3),
    ("Native Storage",
     "Vector attributes stored inline with graph vertices — no external\n"
     "vector database required. Single query joins graph traversal\n"
     "and vector search atomically.",
     BLUE, 8.7, 1.3),
    ("Installed Queries",
     "vectorSearch() function available in compiled GSQL queries.\n"
     "Results piped directly into graph traversal — one hop from\n"
     "vector match to full entity context.",
     GREEN, 0.4, 3.9),
    ("4 Embedding Slots",
     "Transaction: risk_embedding + behaviour_emb (two semantic angles\n"
     "on the same transaction). Merchant: embedding. User: embedding.\n"
     "All 384-dim FLOAT vectors.",
     AMBER, 4.55, 3.9),
    ("REST++ Integration",
     "Installed vector search queries exposed as REST endpoints.\n"
     "Any application can call Top-K search with a single HTTP GET —\n"
     "no client SDK required.",
     RGBColor(0xD0, 0x60, 0xC0), 8.7, 3.9),
]
for title, body, col, lft, top in caps:
    add_rect(sl, lft, top, 4.0, 2.2, fill=RGBColor(0x14, 0x22, 0x38))
    add_rect(sl, lft, top, 0.07, 2.2, fill=col)
    add_text_box(sl, title, lft + 0.2, top + 0.15, 3.7, 0.4,
                 font_size=14, bold=True, color=col)
    add_text_box(sl, body, lft + 0.2, top + 0.6, 3.7, 1.5,
                 font_size=11, color=WHITE, wrap=True)

add_rect(sl, 0.4, 6.3, 12.5, 0.55, fill=RGBColor(0x1A, 0x2A, 0x45))
add_text_box(sl, "Key advantage over standalone vector databases (Pinecone, Weaviate, Qdrant):  "
             "TigerGraph combines the query in one engine — vector Top-K narrows candidates, "
             "then graph traversal expands full context without a second round-trip.",
             0.6, 6.35, 12.1, 0.45, font_size=11, color=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EMBEDDING STRATEGY OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "Embedding Strategy  —  How We Encode Meaning",
             "Four semantic embedding channels derived from raw transaction attributes", bar_color=TEAL)

add_text_box(sl,
    "Each entity is encoded into a 384-dimensional vector using sentence-transformers (all-MiniLM-L6-v2). "
    "The vector captures semantic meaning from a carefully engineered text summary — not raw attribute values.",
    0.4, 1.2, 12.5, 0.5, font_size=13, color=NAVY, wrap=True)

channels = [
    ("Transaction\nRisk Embedding",
     "risk_embedding",
     "Encodes financial crime signals:\namount band, merchant category, card type,\n"
     "geography, time-of-day, risk tags\n(e.g. 'cash-like', 'prepaid', 'late-night').",
     ORANGE, 0.35, 2.0),
    ("Transaction\nBehaviour Embedding",
     "behaviour_emb",
     "Encodes spending pattern:\ncategory intent, channel type, day type,\n"
     "merchant city and region, transaction timing.\n"
     "No risk labels — pure behavioural signal.",
     BLUE, 3.55, 2.0),
    ("Merchant\nEmbedding",
     "Merchant.embedding",
     "Aggregates all transactions at the merchant:\nspend category, geography, dominant card type,\n"
     "acquirer info, typical amount range.\n"
     "Captures merchant risk & behaviour profile.",
     TEAL, 6.75, 2.0),
    ("User\nEmbedding",
     "USER.embedding",
     "Aggregates all transactions by the user:\nspend categories visited, geographies,\n"
     "card types used, time preferences.\n"
     "Captures user behavioural fingerprint.",
     GREEN, 9.95, 2.0),
]

for name, attr, desc, col, lft, top in channels:
    add_rect(sl, lft, top, 3.05, 4.6, fill=NAVY)
    add_rect(sl, lft, top, 3.05, 0.08, fill=col)
    add_text_box(sl, name, lft + 0.12, top + 0.15, 2.8, 0.55,
                 font_size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_rect(sl, lft + 0.5, top + 0.8, 2.05, 0.35, fill=RGBColor(0x1A, 0x2E, 0x3A))
    add_text_box(sl, attr, lft + 0.5, top + 0.83, 2.05, 0.3,
                 font_size=9, color=AMBER, align=PP_ALIGN.CENTER, bold=True)
    add_text_box(sl, desc, lft + 0.12, top + 1.3, 2.8, 2.8,
                 font_size=11, color=WHITE, wrap=True)

add_rect(sl, 0.35, 6.75, 12.6, 0.45, fill=RGBColor(0xEB, 0xF5, 0xFF))
add_text_box(sl,
    "Model: all-MiniLM-L6-v2  |  Output: 384-dim float32 vectors  |  "
    "Similarity: COSINE  |  Index: HNSW  |  Stored inline in TigerGraph vertex attributes",
    0.55, 6.78, 12.2, 0.38, font_size=11, color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRANSACTION RISK EMBEDDING STEPS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Transaction Risk Embedding  —  Step-by-Step Construction",
             "risk_embedding  ·  384 dimensions  ·  HNSW COSINE  ·  Targets financial crime signals",
             bar_color=ORANGE)

steps_risk = [
    ("STEP 1", "Load Raw Transaction",
     "Read tran_sequence_number, Transaction_Amount, merchant_category_code,\n"
     "merchant_city, merchant_country_code, merchant_region_code, card_type,\n"
     "transaction_date, transaction_time, issr_id, acquirer_id."),
    ("STEP 2", "Feature Engineering",
     "• Amount → amount_band ('micro' / 'small-ticket' / 'mid-value' / 'high-value' / 'very high-value')\n"
     "• MCC code → category text ('consumer electronics', 'financial services', …)\n"
     "• card_type → 'debit card' / 'credit card' / 'prepaid card'\n"
     "• transaction_time → time_of_day ('morning' / 'afternoon' / 'evening' / 'late-night')\n"
     "• transaction_date → weekday / weekend flag"),
    ("STEP 3", "Risk Tag Derivation",
     "Apply deterministic rules to generate risk signal tags:\n"
     "• Amount ≥ 2,000 → 'very high-value'\n"
     "• Amount ≥ 500  → 'high-value'\n"
     "• MCC 4511/7011 → 'travel-related'\n"
     "• MCC 6012      → 'cash-like financial services'\n"
     "• MCC 5732      → 'premium electronics'\n"
     "• card_type PREPAID → 'prepaid instrument'\n"
     "• hour < 06:00  → 'late-night activity'"),
    ("STEP 4", "Text Assembly",
     "Concatenate into a single risk-focused sentence:\n"
     "\"This transaction appears to be a {amount_band} {category} payment using a {card_type}.\n"
     " It occurred at {merchant_name} in {city}, {country}, within {region},\n"
     " on {date} during {time_of_day} {day_type} hours. Processed in {currency}\n"
     " through issuer {issr_id} and acquirer {acquirer_id}. Described as: {risk_tags}.\""),
    ("STEP 5", "Embed & Store",
     "• Pass text to SentenceTransformer('all-MiniLM-L6-v2').encode()\n"
     "• Normalise to unit vector (L2 normalisation)\n"
     "• Format: pipe-separated float string for TigerGraph bulk load\n"
     "• Loaded to Transaction vertex attribute: risk_embedding (VECTOR[384])"),
]

y = 1.25
for step_id, step_name, step_body in steps_risk:
    add_rect(sl, 0.35, y, 1.2, 0.85, fill=ORANGE)
    add_text_box(sl, step_id,   0.35, y + 0.06, 1.2, 0.3,  font_size=9,  bold=True,  color=NAVY,  align=PP_ALIGN.CENTER)
    add_text_box(sl, step_name, 0.35, y + 0.35, 1.2, 0.35, font_size=10, bold=False, color=NAVY,  align=PP_ALIGN.CENTER)
    add_rect(sl, 1.6, y, 11.35, 0.85, fill=RGBColor(0x12, 0x22, 0x38))
    add_text_box(sl, step_body, 1.7, y + 0.05, 11.1, 0.75, font_size=10, color=WHITE, wrap=True)
    y += 0.97


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TRANSACTION BEHAVIOUR EMBEDDING STEPS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Transaction Behaviour Embedding  —  Step-by-Step Construction",
             "behaviour_emb  ·  384 dimensions  ·  HNSW COSINE  ·  Captures spending patterns, not risk labels",
             bar_color=BLUE)

steps_beh = [
    ("STEP 1", "Same Source Fields",
     "Uses identical raw transaction CSV fields as the risk embedding — "
     "amount, MCC, card type, city, country, region, date, time."),
    ("STEP 2", "Feature Engineering (shared)",
     "• amount_band — same scale as risk embedding\n"
     "• category_text — MCC → human-readable category\n"
     "• card_type_text — CREDIT / DEBIT / PREPAID → natural language\n"
     "• time_of_day + day_type — temporal semantics"),
    ("STEP 3", "No Risk Tags",
     "Unlike the risk embedding, NO risk tags are appended.\n"
     "The goal is to encode pure spending behaviour — what, where, when, how —\n"
     "without any judgement about fraud likelihood."),
    ("STEP 4", "Text Assembly (behaviour-focused)",
     "\"This transaction was a {amount_band} {card_type} payment of {amount} {currency}\n"
     " at {merchant_name}, a {category} merchant in {city}, {country}, within the\n"
     " {region} region. It took place on {date} during {time_of_day} {day_type} hours.\""),
    ("STEP 5", "Embed & Store",
     "• SentenceTransformer('all-MiniLM-L6-v2').encode() → 384-dim vector\n"
     "• L2-normalised before storage\n"
     "• Loaded to Transaction vertex: behaviour_emb (VECTOR[384])\n"
     "• Two vectors on the same transaction → query against either for different insight"),
]

y = 1.25
for step_id, step_name, step_body in steps_beh:
    add_rect(sl, 0.35, y, 1.2, 0.85, fill=BLUE)
    add_text_box(sl, step_id,   0.35, y + 0.06, 1.2, 0.3,  font_size=9,  bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, step_name, 0.35, y + 0.35, 1.2, 0.35, font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.6, y, 11.35, 0.85, fill=RGBColor(0x12, 0x22, 0x38))
    add_text_box(sl, step_body, 1.7, y + 0.05, 11.1, 0.75, font_size=10, color=WHITE, wrap=True)
    y += 0.97


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MERCHANT EMBEDDING STEPS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Merchant Embedding  —  Step-by-Step Construction",
             "Merchant.embedding  ·  384 dimensions  ·  Aggregates all transactions at the merchant",
             bar_color=TEAL)

steps_merch = [
    ("STEP 1", "Group Transactions",
     "For each unique merchant_id, collect ALL transactions.\n"
     "Derive: total_transactions, mean_amount, dominant_card_type (mode),\n"
     "dominant_currency (mode), dominant_category, city, country, region."),
    ("STEP 2", "Compute Transaction Statistics",
     "• Amount statistics: min, max, mean across all merchant transactions\n"
     "• Card type distribution: most frequent card type\n"
     "• MCC profile: category text for the merchant's assigned MCC\n"
     "• Geography: city + country + region from merchant attributes"),
    ("STEP 3", "Derive Semantic Labels",
     "• Amount profile: 'high average transaction value', 'micro-payment focus', etc.\n"
     "• Risk posture: flagged if dominant risk tags contain 'cash-like', 'prepaid', 'late-night'\n"
     "• Channel profile: 'predominantly card-present', 'heavy online / card-not-present'\n"
     "• Geographic type: 'domestic single-city', 'multi-regional', 'international cross-border'"),
    ("STEP 4", "Text Assembly",
     "\"Merchant {name} in {city}, {country} ({region} region) operates in the "
     "{category} sector.\n"
     " It processes primarily {card_type} transactions with an average value of {mean_amount} {currency},\n"
     " serving customers {geographic_profile}. Merchant risk profile: {risk_tags}.\""),
    ("STEP 5", "Embed & Store",
     "• SentenceTransformer('all-MiniLM-L6-v2').encode() → 384-dim vector\n"
     "• L2-normalised\n"
     "• Loaded to Merchant vertex: embedding (VECTOR[384])\n"
     "• Used by Merchant Similarity search and GraphSAGE neighbourhood aggregation"),
]

y = 1.25
for step_id, step_name, step_body in steps_merch:
    add_rect(sl, 0.35, y, 1.2, 0.85, fill=TEAL)
    add_text_box(sl, step_id,   0.35, y + 0.06, 1.2, 0.3,  font_size=9,  bold=True,  color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(sl, step_name, 0.35, y + 0.35, 1.2, 0.35, font_size=10, bold=False, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.6, y, 11.35, 0.85, fill=RGBColor(0x12, 0x22, 0x38))
    add_text_box(sl, step_body, 1.7, y + 0.05, 11.1, 0.75, font_size=10, color=WHITE, wrap=True)
    y += 0.97


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — USER EMBEDDING STEPS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "User Embedding  —  Step-by-Step Construction",
             "USER.embedding  ·  384 dimensions  ·  Aggregates all transactions by the cardholder",
             bar_color=GREEN)

steps_user = [
    ("STEP 1", "Group Transactions",
     "For each unique userid, collect ALL transactions.\n"
     "Derive: total_transactions, unique_merchants_visited, unique_countries,\n"
     "mean_amount, dominant_card_type, dominant_categories."),
    ("STEP 2", "Compute Behavioural Statistics",
     "• Spend categories: top 3 MCC categories by transaction count\n"
     "• Geographic breadth: domestic-only / regional / international\n"
     "• Time preferences: most active time_of_day, weekday vs weekend preference\n"
     "• Card type: primary card type across all transactions"),
    ("STEP 3", "Derive Semantic Labels",
     "• Spend profile: 'premium international spender', 'everyday domestic shopper', etc.\n"
     "• Loyalty signals: repeat merchant visits, loyalty category\n"
     "• Risk signals: frequency of late-night, prepaid, or cash-like transactions\n"
     "• Velocity: transaction frequency, average interval between transactions"),
    ("STEP 4", "Text Assembly",
     "\"User {userid} is a {spend_profile} cardholder who primarily uses {card_type}.\n"
     " They transact across {n_categories} categories including {top_categories},\n"
     " predominantly in {geographic_profile}. Most active during {time_profile}.\n"
     " Risk characteristics: {risk_description}.\""),
    ("STEP 5", "Embed & Store",
     "• SentenceTransformer('all-MiniLM-L6-v2').encode() → 384-dim vector\n"
     "• L2-normalised\n"
     "• Loaded to USER vertex: embedding (VECTOR[384])\n"
     "• Used by User Similarity search for cohort analysis and fraud ring detection"),
]

y = 1.25
for step_id, step_name, step_body in steps_user:
    add_rect(sl, 0.35, y, 1.2, 0.85, fill=GREEN)
    add_text_box(sl, step_id,   0.35, y + 0.06, 1.2, 0.3,  font_size=9,  bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, step_name, 0.35, y + 0.35, 1.2, 0.35, font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.6, y, 11.35, 0.85, fill=RGBColor(0x12, 0x22, 0x38))
    add_text_box(sl, step_body, 1.7, y + 0.05, 11.1, 0.75, font_size=10, color=WHITE, wrap=True)
    y += 0.97


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — HYBRID SEARCH FLOW
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Hybrid Search  —  How a Query Becomes Results",
             "Natural language → embedding → vector Top-K → graph expansion → executive summary",
             bar_color=BLUE)

flow = [
    ("1\nNL Query",
     "\"Suspicious high-value\nlate-night cash\ntransactions\"",
     ORANGE, 0.3),
    ("2\nEmbed",
     "sentence-transformers\nall-MiniLM-L6-v2\n→ 384-dim vector",
     BLUE, 2.75),
    ("3\nVector Search",
     "TigerGraph HNSW\nTop-K (3–20)\nCOSINE similarity",
     TEAL, 5.2),
    ("4\nGraph Expand",
     "GSQL traversal\nfollows edges to\nmerchants, users,\nlocations, MCC",
     GREEN, 7.65),
    ("5\nResults",
     "Executive summary\nTransactions table\nMerchants · Users\nGraph network viz",
     AMBER, 10.1),
]

for step_text, desc, col, lft in flow:
    add_rect(sl, lft, 1.6, 2.15, 3.2, fill=RGBColor(0x12, 0x22, 0x38))
    add_rect(sl, lft, 1.6, 2.15, 0.08, fill=col)
    add_text_box(sl, step_text, lft + 0.1, 1.7, 1.95, 0.55,
                 font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(sl, desc, lft + 0.12, 2.4, 1.9, 2.2,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    if lft < 10.1:
        add_text_box(sl, "▶", lft + 2.15, 3.1, 0.4, 0.4, font_size=18, color=SILVER,
                     align=PP_ALIGN.CENTER)

# Example query strip
add_rect(sl, 0.3, 5.1, 12.7, 0.55, fill=RGBColor(0x10, 0x2A, 0x45))
add_text_box(sl, "Example query:",
             0.5, 5.15, 1.3, 0.42, font_size=11, color=SILVER)
add_text_box(sl, '"Suspicious high-value late-night cash transactions"',
             1.85, 5.15, 7.0, 0.42, font_size=11, bold=True, color=AMBER)
add_text_box(sl, "→  8 similar transactions found · 5 merchants · graph expanded in <500ms",
             9.0, 5.15, 4.0, 0.42, font_size=11, color=GREEN)

# Mode table
add_text_box(sl, "Search Modes", 0.3, 5.85, 3.0, 0.35, font_size=13, bold=True, color=WHITE)
modes = [
    ("🔍 Transaction Risk",     "risk_embedding",   "Financial crime signals"),
    ("💳 Transaction Behaviour","behaviour_emb",    "Spending patterns"),
    ("🏪 Merchant Similarity",  "Merchant.embedding","Risk/behaviour profile"),
    ("👤 User Similarity",      "USER.embedding",   "Cohort / fraud ring detection"),
]
my = 6.25
for mode_name, vec_attr, purpose in modes:
    add_text_box(sl, mode_name,  0.3,  my, 3.4, 0.28, font_size=11, color=WHITE)
    add_text_box(sl, vec_attr,   3.8,  my, 3.4, 0.28, font_size=11, color=AMBER)
    add_text_box(sl, purpose,    7.3,  my, 5.5, 0.28, font_size=11, color=SILVER)
    my += 0.28


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — USE CASES (2×2 grid)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "Four Hybrid Search Use Cases",
             "Each powered by a dedicated embedding + GSQL context expansion query", bar_color=ORANGE)

ucases = [
    ("🔍  Transaction Risk Search", ORANGE,
     "Detect financial crime patterns without predefined rules.",
     ['"Suspicious high-value late-night cash transactions"',
      '"Cross-border payments with unusual geography"',
      '"Prepaid card activity in financial services merchants"'],
     "Get_transaction_context  +  Summarize_transaction_patterns",
     0.3, 1.3),
    ("💳  Transaction Behaviour", BLUE,
     "Surface spending patterns for cohort and anomaly analysis.",
     ['"Premium international travel and airline spend"',
      '"Weekend dining and entertainment patterns"',
      '"Luxury retail and high-end fashion purchases"'],
     "Get_transaction_context  +  Summarize_transaction_patterns",
     6.85, 1.3),
    ("🏪  Merchant Similarity", TEAL,
     "Find merchants with matching risk posture or spend profile.",
     ['"High-risk electronics merchants, large average transactions"',
      '"Online marketplace merchants, high card-not-present volume"',
      '"Fuel stations with frequent micro-payments"'],
     "Get_merchant_context  +  Summarize_merchant_exposure",
     0.3, 4.4),
    ("👤  User Similarity", GREEN,
     "Identify users with comparable behaviour for fraud rings / cohorts.",
     ['"High-value international spenders, frequent cross-border"',
      '"Users with concentrated late-night cash-like activity"',
      '"Dormant accounts with sudden transaction velocity"'],
     "Get_user_context  +  Summarize_user_behavior",
     6.85, 4.4),
]

for title, col, subtitle, queries, gsql, lft, top in ucases:
    add_rect(sl, lft, top, 6.2, 2.8, fill=NAVY)
    add_rect(sl, lft, top, 6.2, 0.06, fill=col)
    add_text_box(sl, title,    lft + 0.15, top + 0.12, 5.9, 0.4, font_size=14, bold=True, color=col)
    add_text_box(sl, subtitle, lft + 0.15, top + 0.55, 5.9, 0.35, font_size=11, color=SILVER)
    for i, q in enumerate(queries):
        add_text_box(sl, f"▸  {q}", lft + 0.15, top + 0.95 + i * 0.38, 5.9, 0.35,
                     font_size=10, color=WHITE)
    add_rect(sl, lft + 0.1, top + 2.42, 6.0, 0.28, fill=RGBColor(0x1A, 0x2A, 0x40))
    add_text_box(sl, gsql, lft + 0.15, top + 2.45, 5.8, 0.22, font_size=9, color=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — GRAPHSAGE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "GraphSAGE  —  Neighbourhood-Aware Embeddings",
             "Graph Neural Network aggregation: risk propagates through connected entities",
             bar_color=RGBColor(0xA0, 0x60, 0xFF))

# Left column — concept
add_text_box(sl, "The Problem with Isolated Embeddings", 0.35, 1.3, 5.5, 0.4,
             font_size=14, bold=True, color=WHITE)
add_text_box(sl,
    "A text-only embedding represents a merchant from its own description:\n\n"
    "  \"Electronics merchant in Singapore.\"\n\n"
    "It cannot see that 80% of the users who transact there also appear\n"
    "in known high-risk patterns across the network.",
    0.35, 1.8, 5.5, 1.8, font_size=12, color=WHITE, wrap=True)

add_text_box(sl, "The GraphSAGE Solution", 0.35, 3.7, 5.5, 0.4,
             font_size=14, bold=True, color=RGBColor(0xC0, 0x90, 0xFF))
add_text_box(sl,
    "Aggregate embeddings from connected neighbours:\n\n"
    "  h_merchant  =  normalise(\n"
    "      α × own_embedding\n"
    "    + (1-α) × mean( transaction embeddings )\n"
    "  )\n\n"
    "α = 0.5 gives equal weight to self and neighbourhood.\n"
    "Lower α → neighbourhood drives the result.",
    0.35, 4.15, 5.5, 2.8, font_size=12, color=WHITE, wrap=True)

# Right column — hop diagram
add_text_box(sl, "Propagation Hops", 7.0, 1.3, 5.9, 0.4,
             font_size=14, bold=True, color=WHITE)

hops = [
    ("1-Hop", "Merchant ← Transactions",
     "Merchant absorbs signals from its own transactions.\n"
     "A new merchant immediately inherits the risk\n"
     "profile of its first transaction partners.",
     TEAL, 1.55),
    ("2-Hop", "Merchant ← Transactions ← Users",
     "Merchants absorb user-level risk, which already\n"
     "aggregates behaviour across ALL merchants visited.\n"
     "Fraud rings and shared mule accounts surface here.",
     RGBColor(0xA0, 0x60, 0xFF), 3.25),
    ("Cold-Start", "New entity, no history",
     "Inject synthetic merchant connected to high-risk\n"
     "transactions. GraphSAGE immediately clusters it\n"
     "near at-risk peers — before any fraud confirmed.",
     ORANGE, 5.1),
]

for hop_name, hop_subtitle, hop_body, col, top in hops:
    add_rect(sl, 7.0, top, 6.0, 1.6, fill=RGBColor(0x14, 0x22, 0x38))
    add_rect(sl, 7.0, top, 0.07, 1.6, fill=col)
    add_text_box(sl, hop_name,     7.2, top + 0.12, 5.7, 0.35, font_size=13, bold=True, color=col)
    add_text_box(sl, hop_subtitle, 7.2, top + 0.45, 5.7, 0.28, font_size=11, color=SILVER)
    add_text_box(sl, hop_body,     7.2, top + 0.75, 5.7, 0.75, font_size=11, color=WHITE, wrap=True)

add_rect(sl, 0.35, 7.0, 12.6, 0.3, fill=RGBColor(0x1A, 0x2A, 0x45))
add_text_box(sl, "Production path: TigerGraph ML Workbench enables native GNN training and real-time "
             "inference on the live graph — no data export required.",
             0.5, 7.03, 12.3, 0.25, font_size=10, color=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLUSTER EXPLORER
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "Embedding Cluster Explorer  —  Visual Intelligence",
             "UMAP 2D projections of 384-dimensional embeddings — structure with zero labels",
             bar_color=TEAL)

clusters = [
    ("Transaction Risk Clusters", ORANGE, 0.35, 1.3,
     "Each dot = one transaction encoded by risk signals.\n"
     "Nearby dots share risk characteristics: amount band,\n"
     "category type, card type, time of day.\n\n"
     "Clusters visible:\n"
     "• High-value late-night electronics cluster\n"
     "• Everyday micro-payment grocery cluster\n"
     "• Cross-border travel cluster\n"
     "• Cash-like financial services cluster"),
    ("Transaction Behaviour Clusters", BLUE, 4.55, 1.3,
     "Same transactions, different embedding axis.\n"
     "Now clustered by spending intent rather than risk:\n"
     "dining, travel, subscription, retail, healthcare.\n\n"
     "Same transaction may sit in different clusters\n"
     "on risk vs behaviour view — shows dual-axis\n"
     "intelligence from a single data point."),
    ("Merchant Clusters", TEAL, 8.75, 1.3,
     "Each dot = one merchant.\n"
     "Merchants with similar category, geography,\n"
     "and spend profiles cluster together.\n\n"
     "Used by Merchant Similarity search to find\n"
     "comparable risk posture merchants and in\n"
     "GraphSAGE to identify at-risk neighbourhood."),
]

for title, col, lft, top, body in clusters:
    add_rect(sl, lft, top, 4.0, 4.8, fill=NAVY)
    add_rect(sl, lft, top, 4.0, 0.07, fill=col)
    add_text_box(sl, title, lft + 0.15, top + 0.15, 3.7, 0.4, font_size=13, bold=True, color=col)
    add_text_box(sl, body, lft + 0.15, top + 0.65, 3.7, 3.9, font_size=11, color=WHITE, wrap=True)

add_rect(sl, 0.35, 6.3, 12.6, 0.85, fill=RGBColor(0xEB, 0xF5, 0xFF))
add_text_box(sl, "How to read a UMAP cluster chart", 0.55, 6.33, 4.0, 0.3,
             font_size=12, bold=True, color=NAVY)
readings = [
    "Nearby dots  →  semantically similar profiles",
    "Dense cluster  →  strong natural grouping (no labels required)",
    "Isolated dot  →  unusual / anomalous entity",
    "Hover a dot  →  see ID and key attributes",
]
rx = [0.55, 3.55, 6.55, 9.55]
for x, r in zip(rx, readings):
    add_text_box(sl, r, x, 6.68, 2.8, 0.3, font_size=10, color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DETECTION COMPARISON
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Detection Capability Comparison",
             "Transaction Intelligence vs. traditional analytics approaches", bar_color=ORANGE)

headers = ["Scenario", "Rule Engine", "SQL / BI", "Vector Only", "Vector + Graph\n(This System)"]
col_x    = [0.35, 3.05, 5.0, 7.0, 9.2]
col_w    = [2.6, 1.85, 1.85, 2.1, 3.9]

# Header row
for hdr, x, w in zip(headers, col_x, col_w):
    bg = BLUE if hdr.startswith("Vector + Graph") else RGBColor(0x1A, 0x28, 0x40)
    add_rect(sl, x, 1.3, w, 0.5, fill=bg)
    add_text_box(sl, hdr, x + 0.05, 1.33, w - 0.1, 0.44,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows_data = [
    ("Known fraudster → new merchant",                    "❌",  "❌",  "⚠️",  "✅"),
    ("Fraud ring sharing one merchant",                   "⚠️", "❌",  "❌",  "✅"),
    ("Plain-English risk query",                          "❌",  "❌",  "✅",  "✅"),
    ("New card / merchant cold-start detection",          "❌",  "❌",  "❌",  "✅"),
    ("Cross-entity context (txn + merch + user)",         "❌",  "⚠️", "❌",  "✅"),
    ("Finds patterns not yet labelled",                   "❌",  "❌",  "✅",  "✅"),
    ("Adapts to new pattern (no code change)",            "❌",  "❌",  "✅",  "✅"),
    ("Explainable path through the graph",                "⚠️", "⚠️", "❌",  "✅"),
    ("Seconds to query a new scenario",                   "❌",  "⚠️", "✅",  "✅"),
]

row_y = 1.85
for i, (scenario, r1, r2, r3, r4) in enumerate(rows_data):
    bg = RGBColor(0x10, 0x1E, 0x30) if i % 2 == 0 else RGBColor(0x14, 0x24, 0x3C)
    for x, w in zip(col_x, col_w):
        add_rect(sl, x, row_y, w, 0.47, fill=bg)
    add_text_box(sl, scenario, col_x[0] + 0.05, row_y + 0.07, col_w[0] - 0.1, 0.35, font_size=10, color=WHITE)
    for val, x, w in [(r1, col_x[1], col_w[1]), (r2, col_x[2], col_w[2]),
                       (r3, col_x[3], col_w[3])]:
        col_val = GREEN if "✅" in val else (AMBER if "⚠️" in val else RGBColor(0xC0, 0x40, 0x40))
        add_text_box(sl, val, x, row_y + 0.07, w, 0.35, font_size=13, color=col_val, align=PP_ALIGN.CENTER)
    add_text_box(sl, r4, col_x[4], row_y + 0.07, col_w[4], 0.35, font_size=13,
                 color=GREEN, align=PP_ALIGN.CENTER)
    row_y += 0.47


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TECHNOLOGY STACK
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "Technology Stack",
             "Every component chosen for production-readiness and scalability", bar_color=BLUE)

layers_tech = [
    ("Graph Database", BLUE, [
        ("TigerGraph 4.2.2 Cloud", "Native vector index (HNSW · COSINE), GSQL query language, REST++ API"),
        ("GSQL",                   "Context queries: Get_transaction/merchant/user_context · Summary queries"),
        ("Vector Attributes",      "VECTOR[384] FLOAT on Transaction (×2), Merchant (×1), USER (×1)"),
    ]),
    ("Embedding Layer", TEAL, [
        ("sentence-transformers",  "all-MiniLM-L6-v2 — 384-dim, L2-normalised float32 vectors"),
        ("Feature Engineering",    "Python pipeline: amount banding, MCC mapping, risk tags, temporal semantics"),
        ("Bulk Loader",            "Pipe-separated vector strings → TigerGraph LOAD JOB"),
    ]),
    ("Application Layer", ORANGE, [
        ("Streamlit",              "Multi-page Python web application — Hybrid Search, Cluster Explorer, GraphSAGE"),
        ("Custom REST Client",     "requests.Session with JWT Bearer token — TigerGraph Cloud 4.x auth"),
        ("plotly + networkx",      "Interactive graph network and UMAP cluster visualisations"),
    ]),
    ("Analytics Layer", GREEN, [
        ("umap-learn",             "UMAP dimensionality reduction — 384D → 2D for cluster visualisation"),
        ("GraphSAGE Simulation",   "Python neighbourhood aggregation — 1-hop, 2-hop, cold-start injection"),
        ("scikit-learn",           "Logistic Regression classifier for supervised fraud scoring simulation"),
    ]),
]

ty = 1.3
for layer_name, col, items in layers_tech:
    add_rect(sl, 0.35, ty, 12.6, 1.25, fill=NAVY)
    add_rect(sl, 0.35, ty, 0.07, 1.25, fill=col)
    add_text_box(sl, layer_name, 0.5, ty + 0.12, 2.0, 0.4, font_size=13, bold=True, color=col)
    for j, (tech, desc) in enumerate(items):
        x = 2.6 + j * 3.55
        add_text_box(sl, tech, x, ty + 0.1, 3.4, 0.32, font_size=11, bold=True, color=WHITE)
        add_text_box(sl, desc, x, ty + 0.45, 3.4, 0.7, font_size=10, color=SILVER, wrap=True)
    ty += 1.35


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — BUSINESS VALUE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
slide_header(sl, "Business Value  —  What This Enables",
             "Strategic capabilities unlocked by graph + vector intelligence", bar_color=ORANGE)

metrics = [
    ("Seconds", "to query any new risk pattern — no rule writing, no SQL"),
    ("Zero",    "labelled training data required for vector search"),
    ("4 ×",     "entity types searchable from a single natural-language query"),
    ("384",     "semantic dimensions capturing behaviour, risk, geography, time"),
]
for i, (val, label) in enumerate(metrics):
    x = 0.35 + i * 3.25
    add_rect(sl, x, 1.3, 3.0, 1.2, fill=RGBColor(0x10, 0x22, 0x38))
    add_rect(sl, x, 1.3, 3.0, 0.07, fill=ORANGE)
    add_text_box(sl, val, x + 0.1, 1.4, 2.8, 0.55, font_size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, label, x + 0.1, 1.95, 2.8, 0.45, font_size=10, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

values = [
    ("Real-time Risk Discovery",
     "Analysts describe an emerging pattern in plain English. The system "
     "instantly surfaces similar transactions from 5,000+ records — before "
     "rules are written or labels collected.",
     ORANGE, 0.35, 2.8),
    ("Cross-Entity Intelligence",
     "A single query expands through the graph: find a risky transaction → "
     "see its merchant, all connected users, location, and MCC. "
     "Context in one query, not five.",
     TEAL, 4.55, 2.8),
    ("Cold-Start Protection",
     "GraphSAGE propagates risk signals to new merchants and users "
     "from day one — before any fraud is confirmed against them. "
     "Relationship-based risk, not history-based.",
     BLUE, 8.75, 2.8),
    ("Unsupervised Clustering",
     "UMAP cluster explorer reveals natural risk and behaviour groupings "
     "with zero manual labelling. Analysts explore the embedding space "
     "to discover patterns not yet in the playbook.",
     GREEN, 0.35, 5.15),
    ("Investigator Efficiency",
     "Compliance officers query in natural language rather than writing "
     "SQL. Results include full entity context — no second query needed "
     "to understand who, where, and when.",
     AMBER, 4.55, 5.15),
    ("Scalable Architecture",
     "TigerGraph native HNSW scales to billions of transactions. "
     "The same GSQL queries, vector schema, and Streamlit UI work "
     "unchanged at production scale.",
     RGBColor(0xD0, 0x60, 0xC0), 8.75, 5.15),
]

for title, body, col, lft, top in values:
    add_rect(sl, lft, top, 3.85, 2.0, fill=RGBColor(0x12, 0x22, 0x38))
    add_rect(sl, lft, top, 0.07, 2.0, fill=col)
    add_text_box(sl, title, lft + 0.18, top + 0.12, 3.55, 0.35, font_size=12, bold=True, color=col)
    add_text_box(sl, body,  lft + 0.18, top + 0.5,  3.55, 1.4,  font_size=10, color=WHITE, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — NEXT STEPS / PRODUCTION PATH
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF9, 0xFC))
slide_header(sl, "Production Path  —  From Demo to Enterprise Scale",
             "Validated architecture ready for phased deployment on TigerGraph Cloud", bar_color=BLUE)

roadmap = [
    ("Phase 1\nNOW",
     "Demo Validated",
     ["✅  Graph schema deployed on TigerGraph 4.2.2",
      "✅  4 vector embedding channels (384-dim HNSW)",
      "✅  Hybrid search queries installed (GSQL)",
      "✅  Streamlit demo app running",
      "✅  GraphSAGE simulation (1-hop, 2-hop, cold-start)"],
     GREEN, 0.35, 1.3),
    ("Phase 2\nSHORT",
     "Production Data Load",
     ["→  Connect live transaction feed (CDC / Kafka)",
      "→  Replace synthetic with real transaction data",
      "→  Automate embedding refresh pipeline",
      "→  Productionise GSQL queries + REST++ endpoints",
      "→  Add role-based access to Streamlit"],
     BLUE, 3.55, 1.3),
    ("Phase 3\nMEDIUM",
     "Native GNN Training",
     ["→  Label confirmed-fraud transactions",
      "→  Train GraphSAGE on TigerGraph ML Workbench",
      "→  Real-time GNN inference on new transactions",
      "→  Streaming embedding updates (no batch cycle)",
      "→  Integrate with case management system"],
     ORANGE, 6.75, 1.3),
    ("Phase 4\nLONG",
     "Enterprise Intelligence Platform",
     ["→  Multi-graph (cards, accounts, devices, IPs)",
      "→  Real-time fraud alert API (REST++)",
      "→  Explainability: full GSQL path per result",
      "→  Regulatory audit trail in graph",
      "→  Global multi-region TigerGraph deployment"],
     RGBColor(0xA0, 0x60, 0xFF), 9.95, 1.3),
]

for phase, title, steps, col, lft, top in roadmap:
    add_rect(sl, lft, top, 3.1, 5.7, fill=NAVY)
    add_rect(sl, lft, top, 3.1, 0.08, fill=col)
    add_rect(sl, lft, top + 0.1, 3.1, 0.65, fill=RGBColor(0x0A, 0x18, 0x30))
    add_text_box(sl, phase, lft + 0.1, top + 0.12, 2.9, 0.55,
                 font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(sl, title, lft + 0.1, top + 0.78, 2.9, 0.4,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, step in enumerate(steps):
        c = GREEN if step.startswith("✅") else WHITE
        add_text_box(sl, step, lft + 0.15, top + 1.3 + i * 0.7, 2.8, 0.6,
                     font_size=11, color=c, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CLOSING / CALL TO ACTION
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 0, 0.18, 7.5, fill=ORANGE)
add_rect(sl, 0.18, 3.4, 13.15, 0.08, fill=BLUE)

add_text_box(sl, "TRANSACTION INTELLIGENCE", 0.5, 1.0, 12.0, 0.8,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(sl, "Graph + Vector.  One Query.  Full Context.", 0.5, 2.0, 12.0, 0.55,
             font_size=20, color=TEAL, align=PP_ALIGN.LEFT)

summary_points = [
    "Vector search finds semantically similar entities — no rules, no labels, no SQL",
    "Graph traversal expands the full context: transactions, merchants, users, locations",
    "GraphSAGE propagates risk through the network — cold-start protected from day one",
    "Natural language to results in under 500ms on TigerGraph's HNSW index",
    "Architecture validated — production path to TigerGraph Cloud native GNN is clear",
]
for i, pt in enumerate(summary_points):
    add_rect(sl, 0.5, 3.65 + i * 0.55, 0.35, 0.35, fill=ORANGE)
    add_text_box(sl, pt, 1.0, 3.66 + i * 0.55, 11.5, 0.42, font_size=13, color=WHITE)

add_text_box(sl, "Built on  TigerGraph 4.2.2  ·  sentence-transformers  ·  GSQL  ·  Streamlit",
             0.5, 6.9, 12.0, 0.35, font_size=11, color=SILVER, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\n✅  Presentation saved → {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")
slide_titles = [
    "1  Title",
    "2  The Challenge",
    "3  Platform Architecture",
    "4  Graph Data Model",
    "5  TigerGraph Vector Capabilities",
    "6  Embedding Strategy Overview",
    "7  Transaction Risk Embedding — Steps",
    "8  Transaction Behaviour Embedding — Steps",
    "9  Merchant Embedding — Steps",
    "10  User Embedding — Steps",
    "11  Hybrid Search Flow",
    "12  Four Use Cases",
    "13  GraphSAGE",
    "14  Cluster Explorer",
    "15  Detection Comparison",
    "16  Technology Stack",
    "17  Business Value",
    "18  Production Roadmap",
    "19  Closing / Call to Action",
]
for t in slide_titles:
    print(f"   {t}")
